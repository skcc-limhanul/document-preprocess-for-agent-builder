from pptx import Presentation


def parse_content(file_path):
    '''
    Extracts plain text and tables from a PPTX file.
    Parameters:
    * file_path : The path to pdf
    Returns:
    * parsed_text : Parsed contents like plain texts and markdown style tables
    '''
    prs = Presentation(file_path)
    output = ''
    for slide in prs.slides:
        output += '\n'
        elements = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                elements.append({
                    'type': 'text',
                    'top': shape.top,
                    'left': shape.left,
                    'content': shape.text.strip()
                })
            elif shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    table_data\
                        .append([cell.text.strip() for cell in row.cells])
                elements.append({
                    'type': 'table',
                    'top': shape.top,
                    'left': shape.left,
                    'content': table_data
                })
        elements.sort(key=lambda x: (x['top'], x['left']))
        for el in elements:
            if el['type'] == 'text':
                output += el['content'] + '\n'
            elif el['type'] == 'table':
                table = el['content']
                output += '\n| ' + ' | '.join(table[0]) + ' |\n'
                output += '| ' + ' | '.join(['---'] * len(table[0])) + ' |\n'
                for row in table[1:]:
                    output += '| ' + ' | '.join(row) + ' |\n'
    return output
